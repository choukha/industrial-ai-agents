# Install the packages
! pip install pandas numpy matplotlib openpyxl python-docx reportlab python-pptx

"""
Comprehensive Interconnected Synthetic Industrial Dataset Generator

Generates a diverse, interconnected set of industrial data for multiple plant sites,
including various file formats, realistic tag names, linked events (Alarms, Maintenance,
Downtime, Batches), Quality Reports, Environmental Data, Safety Incidents, SOPs,
RCAs (DOCX, PDF, PPTX), and a metadata inventory file.
Designed for testing RAG systems, knowledge graph construction, and operational analysis.
"""

import os
import pandas as pd
import numpy as np
import json
import matplotlib.pyplot as plt
import csv
from PIL import Image, ImageDraw, ImageFont
import random
from datetime import datetime, timedelta
import openpyxl # For Excel
from openpyxl.styles import Font as OpenpyxlFont
from docx import Document as DocxDocument # For Word
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as ReportLabImage, PageBreak
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from pptx import Presentation  # For PowerPoint
from pptx.util import Inches as PptxInches, Pt as PptxPt
import io
import uuid # For generating unique IDs
import copy # To avoid modifying global data unintentionally in loops

# --- Configuration ---
BASE_DIR_NAME = "manganese_alloys_plant_data"
global base_dir
base_dir = os.path.join(os.getcwd(), BASE_DIR_NAME)
os.makedirs(base_dir, exist_ok=True)
PLANT_SITES = [
    {
        "code": "ENP",
        "name": "Eramet Norway Porsgrunn",
        "location": "Porsgrunn",
        "areas": {
            "RMH": "Raw Material Handling", "SMT": "Smelting", "REF": "Refining",
            "CST": "Casting", "CLN": "Gas Cleaning", "UTR": "Utilities & Recovery",
            "LAB": "Laboratory", "MNT": "Maintenance Workshop", "WHS": "Warehouse"
        },
        "facilities": {
            "eaf_count": 2, "refinery": True, "crushing_plant": True,
            "water_gas_cleaning": True, "heat_recovery": True,
            "other_equipment": ["Ladles", "Casting Machines", "Tapping Stations"]
        }
    },
    {
        "code": "ENS",
        "name": "Eramet Norway Sauda",
        "location": "Sauda",
         "areas": {
            "RMH": "Raw Material Handling", "SMT": "Smelting", "REF": "Refining",
            "CLN": "Gas Cleaning", "UTR": "Utilities & Recovery", "LAB": "Laboratory",
            "MNT": "Maintenance Workshop"
        },
        "facilities": {
            "eaf_count": 2, "refinery": True, "crushing_plant": True,
            "water_gas_cleaning": True, "heat_recovery": True,
            "other_equipment": ["Ladles", "Tapping Stations", "Raw Material Silos"]
        }
    },
    {
        "code": "ENK",
        "name": "Eramet Norway Kvinesdal",
        "location": "Kvinesdal",
         "areas": {
            "RMH": "Raw Material Handling", "SMT": "Smelting", "CLN": "Gas Cleaning",
            "UTR": "Utilities & Recovery", "SLG": "Slag Processing", "LAB": "Laboratory",
            "MNT": "Maintenance Workshop"
        },
        "facilities": {
            "eaf_count": 1, "refinery": False, "crushing_plant": True,
            "water_gas_cleaning": True, "heat_recovery": True,
            "other_equipment": ["Furnace Charging System", "Slag Processing Unit"]
        }
    }
]

# Roles and Safety Config
PERSONNEL_ROLES = ["Furnace Operator", "Refinery Operator", "Crusher Operator", "Maintenance Technician", "Electrician", "Lab Analyst", "Shift Supervisor", "Safety Officer", "Logistics Coordinator"]
INCIDENT_TYPES = ["Near Miss", "First Aid Injury", "Medical Treatment Injury", "Lost Time Injury", "Property Damage", "Environmental Spill"]
SEVERITY_LEVELS = ["Minor", "Moderate", "Serious", "Major"]
COMMON_HAZARDS = ["Molten metal splash", "High temperature exposure", "Gas exposure (CO)", "Dust inhalation", "Slips/Trips/Falls", "Moving machinery", "Electrical shock", "Falling objects", "Chemical handling"]
RCA_METHODS = ["5 Whys", "Fishbone (Ishikawa)", "Fault Tree Analysis (Simplified Text)"]

# Global storage for generated IDs to enable linking (reset per plant)
generated_ids = {
    "alarms": [], "maintenance": [], "downtime": [], "batches": [],
    "incidents": [], "rcas": []
}
# Global storage for generated file metadata (accumulates across plants)
metadata_list = []
# Global storage for operational data needed across functions within a plant run
# (Use cautiously, prefer passing data as arguments where possible)
current_plant_data = {}

# --- Helper Functions ---
def ensure_dir(directory):
    """Creates a directory if it doesn't exist."""
    if not os.path.exists(directory):
        os.makedirs(directory)

def generate_uuid():
    """Generates a short unique ID."""
    return str(uuid.uuid4())[:8].upper()

def get_random_id(id_type, plant_code_filter=None):
    """Gets a random ID from the generated list, optionally filtering by plant code."""
    ids = generated_ids.get(id_type, [])
    if plant_code_filter:
        # Assumes IDs have plant code embedded (e.g., MWO-ENP-XXXX)
        filtered_ids = [id_val for id_val in ids if f"-{plant_code_filter}-" in id_val]
        return random.choice(filtered_ids) if filtered_ids else None
    else:
        return random.choice(ids) if ids else None

def add_metadata(file_path, plant_code, category, file_format, description, date_range=None, equipment_tags=None, linked_ids=None):
    """Adds file metadata to the global list."""
    # Ensure path uses forward slashes for consistency, relative to base_dir
    relative_path = os.path.relpath(file_path, base_dir).replace(os.sep, '/')
    metadata_list.append({
        "file_path": relative_path,
        "plant_code": plant_code,
        "category": category,
        "file_format": file_format,
        "description": description,
        "date_range": date_range, # Can be single date string or range string
        "equipment_tags": equipment_tags or [],
        "linked_ids": linked_ids or {} # e.g., {"AlarmID": ["ALM-123"], "MWO_ID": ["MWO-456"]}
    })

# --- Tag Naming ---
def generate_equipment_tag(plant_code, area_code, equip_type, equip_num):
    """Generates a standard equipment tag."""
    return f"{plant_code}-{area_code}-{equip_type}{equip_num}"

def generate_sensor_tag(equipment_tag, measurement_type, instance=1):
    """Generates a standard sensor tag."""
    return f"{equipment_tag}-{measurement_type}-{instance:03d}"

# --- Get Random Plant Elements ---
def get_random_equipment_tag(plant_code):
    """Gets a random equipment tag for the given plant."""
    equip_list = current_plant_data.get("equipment_list", [])
    return random.choice([e['tag'] for e in equip_list]) if equip_list else f"{plant_code}-GEN-EQUIP{random.randint(1,5)}"

def get_random_area_code(plant_code):
    """Gets a random area code for the given plant."""
    plant_info = next((p for p in PLANT_SITES if p["code"] == plant_code), None)
    return random.choice(list(plant_info["areas"].keys())) if plant_info else "GEN"

# --- 1. Structured Data Generation (CSV/Excel) ---
def generate_sensor_data(plant_site_info, days=7, interval_minutes=15):
    """Generates time series sensor data using standard tags."""
    plant_code = plant_site_info["code"]
    category = "Sensor_Data"
    data_dir = os.path.join(base_dir, plant_code, category)
    ensure_dir(data_dir)
    print(f"  Generating Sensor Data for {plant_code}...")

    end_date = datetime.now()
    start_date = end_date - timedelta(days=days)
    date_range_str = f"{start_date.strftime('%Y%m%d')}_to_{end_date.strftime('%Y%m%d')}"
    date_range_pd = pd.date_range(start=start_date, end=end_date, freq=f"{interval_minutes}min")

    # Define parameters with typical measurement types
    parameters = {
        "TE": {"mean": 1500, "std": 50, "unit": "°C", "equip_types": ["EAF", "REF", "BOILER"]}, # Temperature
        "P": {"mean": 105, "std": 10, "unit": "kPa", "equip_types": ["EAF", "BOILER", "GASCLN"]}, # Pressure
        "W_pow": {"mean": 25, "std": 5, "unit": "MW", "equip_types": ["EAF", "CRUSH"]}, # Power
        "I": {"mean": 70, "std": 10, "unit": "kA", "equip_types": ["EAF"]}, # Current
        "W_tap": {"mean": 50, "std": 5, "unit": "tons", "frequency": 4 * 24, "equip_types": ["EAF"]}, # Tap Weight
        "F_o2": {"mean": 500, "std": 50, "unit": "Nm³/h", "equip_types": ["REF"]}, # O2 Flow
        "W_load": {"mean": 100, "std": 20, "unit": "t/h", "equip_types": ["CRUSH"]}, # Crusher Load
        "A_co": {"mean": 5, "std": 2, "unit": "%", "equip_types": ["GASCLN"]} # CO Analysis
    }

    equipment_list = []
    # Generate tags for major equipment
    for area_code, area_name in plant_site_info["areas"].items():
        if area_code == "SMT":
            for i in range(plant_site_info["facilities"]["eaf_count"]):
                equipment_list.append({"tag": generate_equipment_tag(plant_code, area_code, "EAF", i+1), "type": "EAF"})
        elif area_code == "REF" and plant_site_info["facilities"]["refinery"]:
            equipment_list.append({"tag": generate_equipment_tag(plant_code, area_code, "REF", 1), "type": "Refinery"})
        elif area_code == "RMH" and plant_site_info["facilities"]["crushing_plant"]:
            equipment_list.append({"tag": generate_equipment_tag(plant_code, area_code, "CRUSH", 1), "type": "Crusher"})
        elif area_code == "UTR" and plant_site_info["facilities"]["heat_recovery"]:
            equipment_list.append({"tag": generate_equipment_tag(plant_code, area_code, "BOILER", 1), "type": "Boiler"})
        elif area_code == "CLN" and plant_site_info["facilities"]["water_gas_cleaning"]:
             equipment_list.append({"tag": generate_equipment_tag(plant_code, area_code, "GASCLN", 1), "type": "GasCleaner"})
        # Add other potential equipment tags if needed for sensors later

    current_plant_data["equipment_list"] = equipment_list # Store for other functions

    for equipment in equipment_list:
        equip_tag = equipment["tag"]
        equip_type = equipment["type"]
        equip_data = {"timestamp": date_range_pd}
        sensor_tags_for_equip = []

        # Generate data for relevant parameters for this equipment type
        instance_counters = {} # To handle multiple sensors of same type on one equipment
        for param_key, config in parameters.items():
            if equip_type in config["equip_types"]:
                measurement_type = param_key.split('_')[0] # e.g., TE, P, W
                instance = instance_counters.get(measurement_type, 0) + 1
                instance_counters[measurement_type] = instance
                sensor_tag = generate_sensor_tag(equip_tag, measurement_type, instance)
                sensor_tags_for_equip.append(sensor_tag)

                # Simulate data
                if param_key == "W_tap": # Discrete tapping events
                    series_data = np.zeros(len(date_range_pd))
                    num_taps = len(date_range_pd) // config.get("frequency", len(date_range_pd))
                    tap_indices = sorted(random.sample(range(len(date_range_pd)), max(0, num_taps + random.randint(-2,2))))
                    for idx in tap_indices:
                        series_data[idx] = max(0, np.random.normal(config["mean"], config["std"]))
                else: # Continuous data
                    series_data = np.random.normal(config["mean"], config["std"], size=len(date_range_pd))
                    # Add some anomalies that might trigger alarms/maintenance
                    if param_key == "TE" and random.random() < 0.1: # 10% chance of temp spike
                        anomaly_start = random.randint(0, len(series_data) - 6)
                        series_data[anomaly_start:anomaly_start+5] += np.random.uniform(100, 250)
                    elif param_key == "P" and random.random() < 0.05: # 5% chance of pressure drop
                        anomaly_start = random.randint(0, len(series_data) - 6)
                        series_data[anomaly_start:anomaly_start+5] -= np.random.uniform(20, 40)

                series_data = np.maximum(series_data, 0) # Ensure no negative values
                equip_data[sensor_tag] = series_data
                equip_data[f"{sensor_tag}_unit"] = config["unit"]

        df = pd.DataFrame(equip_data)
        base_filename = f"{equip_tag}_sensor_data_{date_range_str}"
        filename_csv = os.path.join(data_dir, f"{base_filename}.csv")
        filename_xlsx = os.path.join(data_dir, f"{base_filename}.xlsx")

        df.to_csv(filename_csv, index=False)
        df.to_excel(filename_xlsx, index=False)

        # Add metadata entry for both files
        desc = f"Time series sensor data for {equip_tag} ({equip_type}). Tags: {', '.join(sensor_tags_for_equip)}"
        add_metadata(filename_csv, plant_code, category, "CSV", desc, date_range_str, [equip_tag])
        add_metadata(filename_xlsx, plant_code, category, "XLSX", desc, date_range_str, [equip_tag])

    print(f"    Sensor data generation for {plant_code} complete.")
    # No return needed, stored in current_plant_data

# --- 2. Semi-structured Data Generation (JSON, Excel) ---
def generate_operational_logs(plant_site_info, num_days=90):
    """Generates interconnected maintenance logs, alarm records, batch reports."""
    plant_code = plant_site_info["code"]
    category = "Operational_Logs"
    data_dir = os.path.join(base_dir, plant_code, category)
    ensure_dir(data_dir)
    print(f"  Generating Operational Logs for {plant_code}...")

    end_date = datetime.now()
    start_date = end_date - timedelta(days=num_days)
    date_range_str = f"{start_date.strftime('%Y%m%d')}_to_{end_date.strftime('%Y%m%d')}"
    month_str = end_date.strftime('%Y%m')

    equipment_tags_plant = [e['tag'] for e in current_plant_data.get("equipment_list", [])]
    if not equipment_tags_plant:
        equipment_tags_plant = [f"{plant_code}-GEN-EQUIP{i+1}" for i in range(3)] # Fallback

    # --- Generate Alarms ---
    alarm_data = []
    alarm_types = ["High Temperature", "Low Pressure", "Power Fluctuation", "Equipment Offline", "High CO Level", "Low Flow", "Vibration High", "Leak Detected"]
    priorities = ["Low", "Medium", "High", "Critical"]
    num_alarms = 150 + random.randint(-30, 30)
    generated_ids["alarms"] = [] # Reset for this plant

    for _ in range(num_alarms):
        equip_tag = random.choice(equipment_tags_plant)
        ts = start_date + timedelta(days=random.randint(0, num_days-1), seconds=random.randint(0, 86400))
        alarm_id = f"ALM-{plant_code}-{generate_uuid()}" # Include plant code
        alarm_type = random.choice(alarm_types)
        priority = random.choice(priorities)
        acknowledged = random.choice([True, False]) if priority != "Critical" else False
        resolved = acknowledged and random.choice([True, False, False])

        alarm_data.append({
            "AlarmID": alarm_id, "Timestamp": ts.strftime("%Y-%m-%d %H:%M:%S"),
            "EquipmentTag": equip_tag, "AlarmType": alarm_type, "Priority": priority,
            "Message": f"{priority} priority: {alarm_type} detected on {equip_tag}.",
            "Acknowledged": acknowledged, "Resolved": resolved,
            "TriggeredDowntimeID": None, "TriggeredMaintenanceID": None
        })
        generated_ids["alarms"].append(alarm_id)
    current_plant_data["alarm_data"] = alarm_data # Store for linking

    # --- Generate Downtime Events (Linked to some Alarms) ---
    downtime_data = []
    downtime_reasons = ["Unplanned Maintenance", "Planned Maintenance", "Equipment Failure", "Operator Error", "No Raw Material", "Utility Outage", "Safety Stop"]
    num_downtime_events = 40 + random.randint(-10, 10)
    generated_ids["downtime"] = [] # Reset

    # More realistic linking: High/Critical unresolved alarms have higher chance of causing downtime
    alarms_needing_downtime = [a for a in alarm_data if a["Priority"] in ["High", "Critical"] and not a["Resolved"] and random.random() < 0.6]
    num_linked_downtime = min(num_downtime_events // 2, len(alarms_needing_downtime))
    linked_alarm_indices = random.sample(range(len(alarms_needing_downtime)), num_linked_downtime)

    for i in range(num_downtime_events):
        downtime_id = f"DWT-{plant_code}-{generate_uuid()}" # Include plant code
        generated_ids["downtime"].append(downtime_id)
        triggering_alarm_id = None
        equip_tag = None
        reason = None
        start_ts = None

        if i < num_linked_downtime:
            # Link to a critical alarm
            alarm_index = linked_alarm_indices[i]
            triggering_alarm = alarms_needing_downtime[alarm_index]
            triggering_alarm_id = triggering_alarm["AlarmID"]
            equip_tag = triggering_alarm["EquipmentTag"]
            reason = random.choice(["Unplanned Maintenance", "Equipment Failure", "Safety Stop"])
            start_ts = datetime.strptime(triggering_alarm["Timestamp"], "%Y-%m-%d %H:%M:%S") + timedelta(minutes=random.randint(5, 60))
            # Update the original alarm entry IN THE STORED LIST
            current_plant_data["alarm_data"][current_plant_data["alarm_data"].index(triggering_alarm)]["TriggeredDowntimeID"] = downtime_id
        else:
            # Generate independent downtime
            equip_tag = random.choice(equipment_tags_plant)
            reason = random.choice(downtime_reasons)
            start_ts = start_date + timedelta(days=random.randint(0, num_days-1), hours=random.randint(0, 23))

        duration_hours = round(random.uniform(0.5, 24) if "Planned" not in reason else random.uniform(4, 72), 1)
        end_ts = start_ts + timedelta(hours=duration_hours)

        downtime_data.append({
            "DowntimeID": downtime_id, "EquipmentTag": equip_tag,
            "StartTime": start_ts.strftime("%Y-%m-%d %H:%M:%S"),
            "EndTime": end_ts.strftime("%Y-%m-%d %H:%M:%S"),
            "DurationHours": duration_hours, "Reason": reason,
            "TriggeringAlarmID": triggering_alarm_id, "RelatedMWO_ID": None
        })
    current_plant_data["downtime_data"] = downtime_data # Store for linking

    # --- Generate Maintenance Logs (Linked to some Alarms/Downtime) ---
    maintenance_data = []
    maintenance_types = ["Preventive", "Corrective", "Predictive", "Breakdown", "Improvement"]
    actions = ["Repair", "Replace Component", "Inspect", "Calibrate", "Clean", "Lubricate", "Install Upgrade", "Troubleshoot"]
    technicians = ["J. Olsen", "M. Hansen", "L. Pedersen", "A. Nilsen", "S. Eriksen", "T. Nguyen"]
    parts = ["None", "Filter XYZ", "Sensor ABC", "Bearing 123", "Electrode Clamp", "Refractory Brick", "Pump Seal", "Valve Actuator", "Control Module"]
    num_maint_entries = 70 + random.randint(-15, 15)
    generated_ids["maintenance"] = [] # Reset

    # Link maintenance to some medium+ alarms and relevant downtime
    alarms_triggering_maint = [a for a in current_plant_data["alarm_data"] if a["Priority"] in ["Medium", "High", "Critical"] and not a["Resolved"] and random.random() < 0.4]
    downtime_triggering_maint = [d for d in current_plant_data["downtime_data"] if ("Maintenance" in d["Reason"] or "Failure" in d["Reason"]) and random.random() < 0.7]
    events_triggering_maint = random.sample(alarms_triggering_maint, k=min(num_maint_entries // 4, len(alarms_triggering_maint))) + \
                              random.sample(downtime_triggering_maint, k=min(num_maint_entries // 3, len(downtime_triggering_maint)))

    linked_event_indices = random.sample(range(len(events_triggering_maint)), len(events_triggering_maint)) # Shuffle

    for i in range(num_maint_entries):
        mwo_id = f"MWO-{plant_code}-{generate_uuid()}" # Include plant code
        generated_ids["maintenance"].append(mwo_id)
        triggering_alarm_id = None
        related_downtime_id = None
        equip_tag = None
        maint_type = None
        description = None
        maint_date = None

        if i < len(linked_event_indices):
            event_index = linked_event_indices[i]
            triggering_event = events_triggering_maint[event_index]
            equip_tag = triggering_event["EquipmentTag"]

            if "AlarmID" in triggering_event:
                triggering_alarm_id = triggering_event["AlarmID"]
                maint_type = random.choice(["Corrective", "Breakdown"])
                maint_date = datetime.strptime(triggering_event["Timestamp"], "%Y-%m-%d %H:%M:%S").date() + timedelta(days=random.randint(0, 2))
                description = f"Address alarm {triggering_alarm_id} ({triggering_event['AlarmType']})."
                # Update original alarm IN STORED LIST
                alarm_original_index = next((idx for idx, a in enumerate(current_plant_data["alarm_data"]) if a["AlarmID"] == triggering_alarm_id), -1)
                if alarm_original_index != -1:
                    current_plant_data["alarm_data"][alarm_original_index]["TriggeredMaintenanceID"] = mwo_id
                    current_plant_data["alarm_data"][alarm_original_index]["Resolved"] = True # Assume maint resolves it

            elif "DowntimeID" in triggering_event:
                related_downtime_id = triggering_event["DowntimeID"]
                maint_type = "Corrective" if "Unplanned" in triggering_event["Reason"] else "Planned Maintenance"
                maint_date = datetime.strptime(triggering_event["StartTime"], "%Y-%m-%d %H:%M:%S").date()
                description = f"Maintenance during downtime {related_downtime_id} ({triggering_event['Reason']})."
                # Update original downtime IN STORED LIST
                downtime_original_index = next((idx for idx, d in enumerate(current_plant_data["downtime_data"]) if d["DowntimeID"] == related_downtime_id), -1)
                if downtime_original_index != -1:
                    current_plant_data["downtime_data"][downtime_original_index]["RelatedMWO_ID"] = mwo_id
        else:
            # Generate independent maintenance (often Preventive)
            equip_tag = random.choice(equipment_tags_plant)
            maint_type = random.choice(["Preventive", "Predictive", "Corrective", "Improvement"])
            maint_date = (start_date + timedelta(days=random.randint(0, num_days-1))).date()
            description = f"{maint_type} check/action on {equip_tag.split('-')[-1]}."

        action = random.choice(actions)
        parts_used = random.choice(parts)
        cost = round(random.uniform(50, 2500) + (500 if parts_used != "None" else 0), 2)
        duration = round(random.uniform(1, 8) if maint_type != "Preventive" else random.uniform(2, 16), 1)

        maintenance_data.append({
            "MWO_ID": mwo_id, "Date": maint_date.strftime("%Y-%m-%d"),
            "EquipmentTag": equip_tag, "MaintenanceType": maint_type,
            "Description": f"{action} performed. {description}",
            "Technician": random.choice(technicians), "DurationHours": duration,
            "PartsUsed": parts_used, "Cost": cost, "Status": "Completed",
            "TriggeringAlarmID": triggering_alarm_id, "RelatedDowntimeID": related_downtime_id
        })
    current_plant_data["maintenance_data"] = maintenance_data # Store for linking

    # Save Maintenance Log (Excel)
    maint_df = pd.DataFrame(maintenance_data)
    maint_file_xlsx = os.path.join(data_dir, f"{plant_code}_MaintenanceLog_{month_str}.xlsx")
    maint_df.to_excel(maint_file_xlsx, index=False)
    add_metadata(maint_file_xlsx, plant_code, category, "XLSX", f"Maintenance Work Order logs for {month_str}.", month_str, list(maint_df['EquipmentTag'].unique()), {"AlarmID": list(maint_df['TriggeringAlarmID'].dropna().unique()), "DowntimeID": list(maint_df['RelatedDowntimeID'].dropna().unique())})

    # Save Alarm Records (JSON and Excel)
    alarm_df = pd.DataFrame(current_plant_data["alarm_data"])
    alarm_file_json = os.path.join(data_dir, f"{plant_code}_AlarmRecords_{date_range_str}.json")
    alarm_file_xlsx = os.path.join(data_dir, f"{plant_code}_AlarmRecords_{date_range_str}.xlsx")
    with open(alarm_file_json, 'w') as f:
        json.dump(current_plant_data["alarm_data"], f, indent=2)
    alarm_df.to_excel(alarm_file_xlsx, index=False)
    add_metadata(alarm_file_json, plant_code, category, "JSON", f"Alarm records log.", date_range_str, list(alarm_df['EquipmentTag'].unique()), {"DowntimeID": list(alarm_df['TriggeredDowntimeID'].dropna().unique()), "MWO_ID": list(alarm_df['TriggeredMaintenanceID'].dropna().unique())})
    add_metadata(alarm_file_xlsx, plant_code, category, "XLSX", f"Alarm records log.", date_range_str, list(alarm_df['EquipmentTag'].unique()), {"DowntimeID": list(alarm_df['TriggeredDowntimeID'].dropna().unique()), "MWO_ID": list(alarm_df['TriggeredMaintenanceID'].dropna().unique())})

    # Save Downtime Logs (Excel)
    downtime_df = pd.DataFrame(current_plant_data["downtime_data"])
    downtime_file_xlsx = os.path.join(data_dir, f"{plant_code}_DowntimeLog_{date_range_str}.xlsx")
    downtime_df.to_excel(downtime_file_xlsx, index=False)
    add_metadata(downtime_file_xlsx, plant_code, "Downtime_Logs", "XLSX", f"Equipment downtime log.", date_range_str, list(downtime_df['EquipmentTag'].unique()), {"AlarmID": list(downtime_df['TriggeringAlarmID'].dropna().unique()), "MWO_ID": list(downtime_df['RelatedMWO_ID'].dropna().unique())})


    # --- Generate Batch Reports ---
    batch_report_data = []
    alloy_types = ["FeMn-HC", "SiMn-Std", "FeMn-MC"]
    num_batches = 60 + random.randint(-10, 10)
    generated_ids["batches"] = [] # Reset

    eaf_tags = [e["tag"] for e in current_plant_data.get("equipment_list", []) if e["type"] == "EAF"]
    if not eaf_tags: eaf_tags = [f"{plant_code}-SMT-EAF{i+1}" for i in range(plant_site_info["facilities"]["eaf_count"])] # Fallback

    for i in range(num_batches):
        batch_id = f"BCH-{plant_code}-{generate_uuid()}" # Include plant code
        generated_ids["batches"].append(batch_id)
        batch_date = start_date + timedelta(days=random.randint(0, num_days-1))
        alloy = random.choice(alloy_types)
        eaf_tag = random.choice(eaf_tags)

        raw_mats = {
            "ManganeseOre_kg": round(random.uniform(8000, 12000),0),
            "Coke_kg": round(random.uniform(2000, 3000),0),
            "Quartz_kg": round(random.uniform(500, 1500),0) if "SiMn" in alloy else 0,
            "FeScrap_kg": round(random.uniform(1000, 2000),0) if "FeMn" in alloy else 0,
            "Flux_CaO_kg": round(random.uniform(200, 500), 0) # Added flux
        }
        # Simulate some variability
        energy_kwh = round(random.uniform(25000, 40000) * random.uniform(0.95, 1.05), 0)
        tap_weight_tons = round(sum(raw_mats.values()) * random.uniform(0.5, 0.7) / 1000 * random.uniform(0.9, 1.1), 2)

        # Placeholder quality - will be refined by Lab reports
        mn_pct = round(random.uniform(65, 78), 1)
        si_pct = round(random.uniform(14, 20), 1) if "SiMn" in alloy else round(random.uniform(0.5, 1.5), 1)
        c_pct = round(random.uniform(6, 7.5), 1) if "HC" in alloy else round(random.uniform(1, 2), 1)

        batch_report_data.append({
            "BatchID": batch_id, "Date": batch_date.strftime("%Y-%m-%d"),
            "AlloyType": alloy, "FurnaceTag": eaf_tag,
            **raw_mats,
            "Energy_kWh": energy_kwh, "TapWeight_tons": tap_weight_tons,
            "Mn_content_pct_Target": mn_pct, # Target/Estimate
            "Si_content_pct_Target": si_pct, "C_content_pct_Target": c_pct,
            "Operator": random.choice(PERSONNEL_ROLES) # Use defined roles
        })
    current_plant_data["batch_data"] = batch_report_data # Store for linking

    batch_df = pd.DataFrame(batch_report_data)
    batch_file_xlsx = os.path.join(data_dir, f"{plant_code}_BatchReports_{month_str}.xlsx")
    batch_df.to_excel(batch_file_xlsx, index=False)
    add_metadata(batch_file_xlsx, plant_code, category, "XLSX", f"Production batch reports for {month_str}.", month_str, list(batch_df['FurnaceTag'].unique()))

    print(f"    Operational logs generation for {plant_code} complete.")

# --- 3. Lab Analysis / Quality Reports ---
def generate_lab_reports(plant_site_info, num_days=90):
    """Generates lab analysis reports linked to production batches."""
    plant_code = plant_site_info["code"]
    category = "Lab_Analysis"
    data_dir = os.path.join(base_dir, plant_code, category)
    ensure_dir(data_dir)
    print(f"  Generating Lab Analysis Reports for {plant_code}...")

    end_date = datetime.now()
    start_date = end_date - timedelta(days=num_days)
    month_str = end_date.strftime('%Y%m')

    lab_report_data = []
    analysis_elements = ["Mn", "Si", "C", "P", "S", "Fe"] # Key elements for FeMn/SiMn
    # Get batches generated for this plant
    plant_batches = generated_ids.get("batches", [])
    num_reports = len(plant_batches) # Assume one lab report per batch

    if not plant_batches:
        print(f"    WARN: No batches found for {plant_code} to generate lab reports.")
        return

    # Find batch details to infer alloy type for better results
    batch_details_map = {b["BatchID"]: b for b in current_plant_data.get("batch_data", [])}

    for i in range(num_reports):
        batch_id = plant_batches[i]
        batch_info = batch_details_map.get(batch_id, {})
        alloy_type = batch_info.get("AlloyType", random.choice(["FeMn-HC", "SiMn-Std"])) # Get type or guess

        analysis_date = start_date + timedelta(days=random.randint(i * num_days // num_reports, (i+1) * num_days // num_reports -1)) # Spread analysis dates
        lab_id = f"LAB-{plant_code}-{analysis_date.strftime('%Y%m%d')}-{i+1:03d}"

        # Simulate analysis results based on alloy type
        results = {}
        is_simn = "SiMn" in alloy_type
        is_hc = "HC" in alloy_type

        results["Mn"] = round(random.uniform(64, 79) * random.uniform(0.98, 1.02), 2)
        results["Si"] = round(random.uniform(13, 21) * random.uniform(0.95, 1.05), 2) if is_simn else round(random.uniform(0.4, 1.6) * random.uniform(0.8, 1.2), 2)
        results["C"] = round(random.uniform(6.0, 7.8) * random.uniform(0.98, 1.02), 2) if is_hc else round(random.uniform(1.0, 2.5) * random.uniform(0.9, 1.1), 2)
        results["P"] = round(random.uniform(0.05, 0.35) * random.uniform(0.8, 1.2), 3)
        results["S"] = round(random.uniform(0.005, 0.05) * random.uniform(0.7, 1.3), 3)
        results["Fe"] = max(0, round(100 - results["Mn"] - results["Si"] - results["C"] - results["P"] - results["S"], 2)) # Ensure non-negative

        lab_report_data.append({
            "LabAnalysisID": lab_id, "AnalysisDate": analysis_date.strftime("%Y-%m-%d"),
            "ProductionBatchID": batch_id,
            **{f"{el}_pct": res for el, res in results.items()},
            "Analyst": random.choice([r for r in PERSONNEL_ROLES if "Analyst" in r or "Supervisor" in r])
        })
    current_plant_data["lab_data"] = lab_report_data # Store for use

    lab_df = pd.DataFrame(lab_report_data)
    lab_file_xlsx = os.path.join(data_dir, f"{plant_code}_LabAnalysis_{month_str}.xlsx")
    lab_df.to_excel(lab_file_xlsx, index=False)
    add_metadata(lab_file_xlsx, plant_code, category, "XLSX", f"Lab analysis results for production batches in {month_str}.", month_str, linked_ids={"ProductionBatchID": list(lab_df["ProductionBatchID"].unique())})

    # Optional: Generate individual PDF reports for a few batches
    num_pdf_reports = min(5, len(lab_report_data))
    for i in random.sample(range(len(lab_report_data)), num_pdf_reports): # Sample randomly
        report_item = lab_report_data[i]
        pdf_filename = os.path.join(data_dir, f"{report_item['ProductionBatchID']}_LabReport_{report_item['AnalysisDate']}.pdf")
        doc = SimpleDocTemplate(pdf_filename, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        story.append(Paragraph(f"Certificate of Analysis - {plant_site_info['name']}", styles['h1']))
        story.append(Spacer(1, 12))
        story.append(Paragraph(f"Lab Analysis ID: {report_item['LabAnalysisID']}", styles['Normal']))
        story.append(Paragraph(f"Production Batch ID: {report_item['ProductionBatchID']}", styles['Normal']))
        story.append(Paragraph(f"Analysis Date: {report_item['AnalysisDate']}", styles['Normal']))
        story.append(Paragraph(f"Analyst: {report_item['Analyst']}", styles['Normal']))
        story.append(Spacer(1, 12))
        story.append(Paragraph("Chemical Composition (% by weight):", styles['h3']))

        table_data = [['Element', 'Result (%)']]
        for el in analysis_elements:
            table_data.append([el, f"{report_item.get(f'{el}_pct', 'N/A'):.3f}"])

        table = Table(table_data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(table)
        story.append(Spacer(1, 24))
        story.append(Paragraph("--- End of Report ---", styles['Normal']))

        try:
            doc.build(story)
            add_metadata(pdf_filename, plant_code, category, "PDF", f"Certificate of Analysis for Batch {report_item['ProductionBatchID']}.", report_item['AnalysisDate'], linked_ids={"ProductionBatchID": [report_item['ProductionBatchID']]})
        except Exception as e:
            print(f"    ERROR generating PDF lab report {pdf_filename}: {e}")

    print(f"    Lab Analysis report generation for {plant_code} complete.")

# --- 4. Unstructured Data Generation (DOCX, PDF) ---
def generate_production_report(plant_site_info, num_days=30):
    """Generates Monthly Production Report (DOCX)"""
    plant_code = plant_site_info["code"]
    plant_name = plant_site_info["name"]
    category = "Reports_SOPs" # Keep in same category as SOPs
    data_dir = os.path.join(base_dir, plant_code, category)
    ensure_dir(data_dir)
    print(f"    Generating Production Report for {plant_code}...")

    current_month_year = datetime.now().strftime("%B %Y")
    current_date = datetime.now().strftime("%Y-%m-%d")
    month_str = datetime.now().strftime('%Y_%m')
    start_date_month = datetime.now().replace(day=1) - timedelta(days=1) # Approx start of prev month
    start_date_month = start_date_month.replace(day=1)

    # --- Production Report (DOCX) ---
    doc = DocxDocument()
    doc.add_heading(f'{plant_name} - Monthly Production Report - {current_month_year}', level=1)
    doc.add_paragraph(f"Report Date: {current_date}")
    doc.add_paragraph(f"Plant Code: {plant_code}")
    doc.add_heading('1. Executive Summary', level=2)

    # Summarize data from current_plant_data for the relevant period (approximated)
    batches_this_month = [b for b in current_plant_data.get("batch_data", []) if datetime.strptime(b['Date'], "%Y-%m-%d") >= start_date_month]
    downtime_this_month = [d for d in current_plant_data.get("downtime_data", []) if datetime.strptime(d['StartTime'], "%Y-%m-%d %H:%M:%S") >= start_date_month]
    maint_this_month = [m for m in current_plant_data.get("maintenance_data", []) if datetime.strptime(m['Date'], "%Y-%m-%d") >= start_date_month]

    total_batches = len(batches_this_month)
    total_downtime_hours = sum(d['DurationHours'] for d in downtime_this_month)
    total_maint_cost = sum(m['Cost'] for m in maint_this_month)
    avg_energy_kwh_per_ton = np.mean([b['Energy_kWh'] / b['TapWeight_tons'] for b in batches_this_month if b['TapWeight_tons'] > 0]) if batches_this_month else 0

    summary_text = f"Production for {current_month_year} yielded {total_batches} batches. "
    summary_text += f"Key performance indicators: Avg Energy Consumption ~{avg_energy_kwh_per_ton:.0f} kWh/ton. "
    summary_text += f"Total recorded equipment downtime was {total_downtime_hours:.1f} hours. "
    summary_text += f"Maintenance activities cost approx ${total_maint_cost:,.2f}. "
    key_issue_maint_id = get_random_id("maintenance", plant_code) # Get MWO specific to this plant
    key_issue_alarm_id = get_random_id("alarms", plant_code)
    if key_issue_maint_id:
         summary_text += f"Significant maintenance included MWO {key_issue_maint_id.split('-')[-1]}. " # Show short ID
    if key_issue_alarm_id:
        summary_text += f"Notable alarms included {key_issue_alarm_id.split('-')[-1]}. "
    doc.add_paragraph(summary_text)

    doc.add_heading('2. Production Figures', level=2)
    # Aggregate batch data by alloy type
    prod_summary = {}
    for b in batches_this_month:
        alloy = b['AlloyType']
        if alloy not in prod_summary:
            prod_summary[alloy] = {"count": 0, "tons": 0, "energy": 0}
        prod_summary[alloy]["count"] += 1
        prod_summary[alloy]["tons"] += b['TapWeight_tons']
        prod_summary[alloy]["energy"] += b['Energy_kWh']

    prod_table_data = [("Alloy", "Batches", "Actual (tons)", "Avg Energy (kWh/ton)")]
    for alloy, data in prod_summary.items():
        avg_energy = (data['energy'] / data['tons']) if data['tons'] > 0 else 0
        prod_table_data.append((alloy, str(data['count']), f"{data['tons']:.1f}", f"{avg_energy:.0f}"))

    if len(prod_table_data) > 1:
        table = doc.add_table(rows=1, cols=len(prod_table_data[0]))
        table.style = 'Table Grid'
        hdr_cells = table.rows[0].cells
        for i, header in enumerate(prod_table_data[0]): hdr_cells[i].text = header
        for item in prod_table_data[1:]:
            row_cells = table.add_row().cells
            for i, data in enumerate(item): row_cells[i].text = data
    else:
        doc.add_paragraph("No production batches recorded for this period.")

    doc.add_heading('3. Major Downtime Events (>4h)', level=2)
    significant_downtime = sorted([d for d in downtime_this_month if d['DurationHours'] > 4], key=lambda x: x['DurationHours'], reverse=True)[:5] # Top 5
    if significant_downtime:
        for dt in significant_downtime:
            link_info = f"(Alarm: {dt['TriggeringAlarmID'].split('-')[-1]})" if dt.get('TriggeringAlarmID') else ""
            link_info += f" (MWO: {dt['RelatedMWO_ID'].split('-')[-1]})" if dt.get('RelatedMWO_ID') else ""
            doc.add_paragraph(f"- ID: {dt['DowntimeID'].split('-')[-1]}, Equip: {dt['EquipmentTag']}, Dur: {dt['DurationHours']}h, Reason: {dt['Reason']} {link_info}", style='List Bullet')
    else:
        doc.add_paragraph("No significant downtime events (>4 hours) recorded.")

    doc.add_heading('4. Key Maintenance Activities (Cost > $1k or Breakdown)', level=2)
    significant_maint = sorted([m for m in maint_this_month if (m['Cost'] > 1000 or m['MaintenanceType'] == 'Breakdown')], key=lambda x: x['Cost'], reverse=True)[:5] # Top 5
    if significant_maint:
         for mwo in significant_maint:
             link_info = f"(Alarm: {mwo['TriggeringAlarmID'].split('-')[-1]})" if mwo.get('TriggeringAlarmID') else ""
             link_info += f" (Downtime: {mwo['RelatedDowntimeID'].split('-')[-1]})" if mwo.get('RelatedDowntimeID') else ""
             doc.add_paragraph(f"- MWO: {mwo['MWO_ID'].split('-')[-1]}, Equip: {mwo['EquipmentTag']}, Type: {mwo['MaintenanceType']}, Cost: ${mwo['Cost']:.0f} {link_info}", style='List Bullet')
    else:
        doc.add_paragraph("No high-cost or breakdown maintenance recorded.")

    # (Financial Summary Section - similar to v3)
    doc.add_heading('5. Financial Summary (Simplified)', level=2)
    energy_cost = avg_energy_kwh_per_ton * sum(b['TapWeight_tons'] for b in batches_this_month) * random.uniform(0.05, 0.15) # Simplified energy cost $
    doc.add_paragraph(f"- Estimated Energy Cost: ${energy_cost:,.0f}")
    doc.add_paragraph(f"- Total Maintenance Cost: ${total_maint_cost:,.2f}")
    doc.add_paragraph(f"- Raw Material Cost (Est.): ${random.randint(500000, 1000000):,.2f}")


    doc_path = os.path.join(data_dir, f"{plant_code}_Production_Report_{month_str}.docx")
    doc.save(doc_path)
    add_metadata(doc_path, plant_code, category, "DOCX", f"Monthly production report for {current_month_year}.", month_str, linked_ids={"BatchID": [b['BatchID'] for b in batches_this_month], "DowntimeID": [d['DowntimeID'] for d in significant_downtime], "MWO_ID": [m['MWO_ID'] for m in significant_maint]})
    print(f"    Production Report generation for {plant_code} complete.")



# --- 5. Environmental Reports (PDF) ---
def generate_environmental_reports(plant_site_info, num_months=3):
    """Generates monthly environmental summary reports."""
    plant_code = plant_site_info["code"]
    plant_name = plant_site_info["name"]
    category = "Environmental_Data"
    data_dir = os.path.join(base_dir, plant_code, category)
    ensure_dir(data_dir)
    print(f"  Generating Environmental Reports for {plant_code}...")

    end_date = datetime.now()

    for i in range(num_months):
        report_date = end_date - timedelta(days=i * 30)
        report_month_str = report_date.strftime('%Y_%m')
        report_month_year = report_date.strftime('%B %Y')

        pdf_filename = os.path.join(data_dir, f"{plant_code}_Environmental_Report_{report_month_str}.pdf")
        doc = SimpleDocTemplate(pdf_filename, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        story.append(Paragraph(f"Monthly Environmental Report - {plant_name}", styles['h1']))
        story.append(Paragraph(f"Period: {report_month_year}", styles['h2']))
        story.append(Paragraph(f"Report Date: {end_date.strftime('%Y-%m-%d')}", styles['Normal']))
        story.append(Spacer(1, 12))

        story.append(Paragraph("1. Emissions Summary (Estimated)", styles['h3']))
        emissions_data = [
            ['Parameter', 'Value', 'Unit', 'Limit', 'Compliance'],
            ['CO2 eq.', f"{random.randint(5000, 15000):,}", 'tons/month', 'N/A', 'N/A'],
            ['SOx', f"{random.uniform(10, 50):.1f}", 'tons/month', '60', 'Compliant'],
            ['NOx', f"{random.uniform(5, 25):.1f}", 'tons/month', '30', 'Compliant'],
            ['Dust (Stack)', f"{random.uniform(0.5, 5):.2f}", 'tons/month', '5.0', 'Compliant' if random.random() > 0.05 else 'Action Required'],
            ['Dust (Fugitive)', f"{random.uniform(1, 10):.1f}", 'tons/month', 'N/A', 'Monitored'],
        ]
        # Add reference to gas cleaning equipment
        gas_cleaner_tag = next((e['tag'] for e in current_plant_data.get("equipment_list", []) if e['type'] == 'GasCleaner'), None)
        if gas_cleaner_tag:
            story.append(Paragraph(f"(Stack emissions monitored via sensors on {gas_cleaner_tag})", styles['Italic']))


        table = Table(emissions_data, colWidths=[100, 80, 80, 80, 80])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.darkgreen),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.lightgrey),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ('ALIGN', (1, 1), (1, -1), 'RIGHT'), # Align numbers right
            ('ALIGN', (3, 1), (3, -1), 'RIGHT'),
        ]))
        story.append(table)
        story.append(Spacer(1, 12))
        story.append(Paragraph("2. Water Usage & Discharge", styles['h3']))
        # Change: Use styles['Normal'] and add bullet manually
        story.append(Paragraph(f"• Total Water Intake: {random.randint(10000, 50000):,} m³/month", styles['Normal']))
        story.append(Paragraph(f"• Discharge Volume: {random.randint(8000, 45000):,} m³/month", styles['Normal']))
        story.append(Paragraph(f"• Discharge Quality: Meets permit requirements (Ref: Permit {plant_code}-ENV-WTR-001)", styles['Normal']))
        story.append(Spacer(1, 12))

        story.append(Paragraph("3. Waste Management", styles['h3']))
        slag_tag = next((e['tag'] for e in current_plant_data.get("equipment_list", []) if e['type'] == 'Slag Processing Unit'), f"{plant_code}-SLG-PROC1" if "SLG" in plant_site_info["areas"] else None)
        slag_text = f"(Processed via {slag_tag})" if slag_tag else "(Partially reused/sold)"

        # Change: Create a ListBullet style
        list_bullet_style = copy.deepcopy(styles['Normal']) # Ensure new instance
        list_bullet_style.spaceBefore = 6 # Add a bit of space
        list_bullet_style.leftIndent = 24 # Indent for bullet point

        # Apply style directly
        story.append(Paragraph(f"- Slag Produced: {random.randint(1500, 4000):,} tons/month {slag_text}", style=list_bullet_style))
        story.append(Paragraph(f"- Hazardous Waste Disposed: {random.uniform(1, 5):.1f} tons/month (Manifest ID: HW-{plant_code}-{report_month_str})", style=list_bullet_style))
        story.append(Paragraph(f"- Recycled Materials: {random.uniform(10, 50):.1f} tons/month (Cardboard, Scrap Metal)", style=list_bullet_style))

        try:
            doc.build(story)
            add_metadata(pdf_filename, plant_code, category, "PDF", f"Monthly environmental summary report for {report_month_year}.", report_month_str, equipment_tags=[gas_cleaner_tag] if gas_cleaner_tag else [])
        except Exception as e:
            print(f"    ERROR generating PDF Env Report {pdf_filename}: {e}")

    print(f"    Environmental report generation for {plant_code} complete.")

# --- 6. Rich Media Generation (PNG, PDF for Diagrams) ---
def create_diagram_pil(filename, title_text, content_lines, plant_code="", category="Diagrams_Visuals"):
    """Creates a simple diagram using PIL and saves as PNG."""
    diagram_dir = os.path.join(base_dir, plant_code, category)
    ensure_dir(diagram_dir)
    filepath = os.path.join(diagram_dir, filename)

    img_width = 800
    img_height = 600
    bg_color = (230, 230, 230)
    text_color = (0, 0, 0)
    try:
        # Use a common font if available, otherwise default
        # Check for Arial on common OS paths or use default
        font_path = None
        if os.path.exists("/usr/share/fonts/truetype/msttcorefonts/arial.ttf"): # Linux
             font_path = "/usr/share/fonts/truetype/msttcorefonts/arial.ttf"
        elif os.path.exists("/Library/Fonts/Arial.ttf"): # MacOS
             font_path = "/Library/Fonts/Arial.ttf"
        elif os.path.exists("C:\\Windows\\Fonts\\arial.ttf"): # Windows
             font_path = "C:\\Windows\\Fonts\\arial.ttf"

        font_title = ImageFont.truetype(font_path, 24) if font_path else ImageFont.load_default()
        font_content = ImageFont.truetype(font_path, 16) if font_path else ImageFont.load_default()
    except IOError:
        print("WARN: Arial font not found, using default PIL font.")
        font_title = ImageFont.load_default()
        font_content = ImageFont.load_default()
    except ImportError:
         print("WARN: Pillow font module not fully available, using default PIL font.")
         font_title = ImageFont.load_default()
         font_content = ImageFont.load_default()


    img = Image.new('RGB', (img_width, img_height), color=bg_color)
    draw = ImageDraw.Draw(img)

    # Title
    try:
        # textbbox requires Pillow 9.0.0+ and a TrueType font object
        if hasattr(draw, 'textbbox') and isinstance(font_title, ImageFont.FreeTypeFont):
             title_bbox = draw.textbbox((0, 0), title_text, font=font_title)
             title_width = title_bbox[2] - title_bbox[0]
             draw.text(((img_width - title_width) / 2, 20), title_text, font=font_title, fill=text_color)
        else:
             # Fallback for default fonts or older Pillow versions
             title_width, title_height = draw.textsize(title_text, font=font_title)
             draw.text(((img_width - title_width) / 2, 20), title_text, font=font_title, fill=text_color)
    except Exception as e: # Catch potential errors with text rendering
         print(f"WARN: Error rendering title text: {e}. Using fallback.")
         title_width, title_height = draw.textsize(title_text, font=font_title)
         draw.text(((img_width - title_width) / 2, 20), title_text, font=font_title, fill=text_color)


    # Content
    y_text = 80
    for line in content_lines:
        try:
            draw.text((30, y_text), line, font=font_content, fill=text_color)
            y_text += 25
        except Exception as e:
            print(f"WARN: Error rendering content line '{line}': {e}. Skipping line.")
        if y_text > img_height - 40: # Basic overflow prevention
            break

    # Add a border for diagrammatic look
    draw.rectangle([(10,10), (img_width-10, img_height-10)], outline=(50,50,50), width=2)

    img.save(filepath)
    return filepath

def generate_simplified_pid(filename_pdf, plant_code):
    """Generates a simplified P&ID excerpt as a PDF using ReportLab & Matplotlib."""
    category = "Diagrams_Visuals"
    diagram_dir = os.path.join(base_dir, plant_code, category)
    ensure_dir(diagram_dir)
    filepath = os.path.join(diagram_dir, filename_pdf)

    doc = SimpleDocTemplate(filepath, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph(f"P&ID Excerpt - {plant_code} - Raw Material to EAF", styles['h1']))
    story.append(Spacer(1, 12))

    # Find relevant tags from current plant data
    equipment_list = current_plant_data.get("equipment_list", [])
    crusher_tag = next((e['tag'] for e in equipment_list if e['type'] == 'Crusher'), f"{plant_code}-RMH-CRUSH1")
    eaf_tag = next((e['tag'] for e in equipment_list if e['type'] == 'EAF'), f"{plant_code}-SMT-EAF1")
    silo_tag = f"{plant_code}-RMH-SILO1" # Assumed
    conveyor_tag = f"{plant_code}-RMH-CONV1" # Assumed
    weight_sensor_tag = generate_sensor_tag(conveyor_tag, 'W', 1)

    # Create a simple diagram using matplotlib
    try:
        fig, ax = plt.subplots(figsize=(7, 5))
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 7)
        ax.set_xticks([])
        ax.set_yticks([])
        ax.set_title(f"Simplified PFD: {plant_code} Raw Material Input")

        # Silo
        ax.add_patch(plt.Rectangle((1, 3), 2, 3, edgecolor='black', facecolor='lightgray'))
        ax.text(2, 4.5, f'{silo_tag}\n(Ore)', ha='center', va='center', fontsize=8)
        # Crusher
        ax.add_patch(plt.Polygon([[3.5, 4.5], [4.5, 5.5], [5.5, 4.5], [4.5, 3.5]], edgecolor='black', facecolor='lightblue'))
        ax.text(4.5, 4.5, f'{crusher_tag}', ha='center', va='center', fontsize=8)
        # Conveyor
        ax.plot([2, 3.5], [3.5, 4.5], 'k-', lw=2) # Silo to Crusher
        ax.plot([5.5, 7], [4.5, 4.5], 'k-', lw=2) # Crusher to EAF
        ax.text(6.25, 4.7, conveyor_tag, ha='center', va='bottom', fontsize=7)
        # EAF (representation)
        ax.add_patch(plt.Circle((8, 4.5), 1, edgecolor='black', facecolor='orange'))
        ax.text(8, 4.5, f'{eaf_tag}', ha='center', va='center', fontsize=8)
        # Sensor example
        ax.add_patch(plt.Circle((6.5, 4.8), 0.2, edgecolor='green', facecolor='white'))
        ax.text(6.5, 5.2, weight_sensor_tag, ha='center', va='center', fontsize=7) # Weight Transmitter

        img_buffer = io.BytesIO()
        plt.savefig(img_buffer, format='png', bbox_inches='tight')
        plt.close(fig)
        img_buffer.seek(0)

        story.append(ReportLabImage(img_buffer, width=450, height=320))
        # --- *** FIX: Use 'Italic' style instead of 'Caption' *** ---
        story.append(Paragraph(f"Figure 1: Simplified Process Flow Diagram for {silo_tag} -> {crusher_tag} -> {eaf_tag}.", styles['Italic']))

    except Exception as e:
        print(f"    ERROR generating Matplotlib diagram for P&ID: {e}")
        story.append(Paragraph(f"Error generating diagram: {e}", styles['Italic']))


    try:
        doc.build(story)
        add_metadata(filepath, plant_code, category, "PDF", f"Simplified P&ID/PFD for Raw Material Handling.", datetime.now().strftime('%Y%m%d'), [silo_tag, crusher_tag, eaf_tag, conveyor_tag, weight_sensor_tag])
    except Exception as e:
        print(f"    ERROR generating PDF P&ID {filepath}: {e}")

    return filepath


def generate_hmi_snapshot(filename_png, plant_code, equipment_tag, equipment_type):
    """Generates a simplified HMI snapshot as PNG using standard tags."""
    category = "Diagrams_Visuals"
    title = f"HMI: {equipment_tag} Status - {datetime.now().strftime('%H:%M:%S')}"
    timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')

    # Simulate some data based on equipment type
    content = [
        f"Equipment: {equipment_tag} ({equipment_type})",
        f"Timestamp: {timestamp}",
        "--- Key Parameters ---"
    ]
    status = random.choice(["RUNNING", "IDLE", "ALARM", "MAINTENANCE"])
    active_alarm_id = None

    if equipment_type == "EAF":
        temp_tag = generate_sensor_tag(equipment_tag, 'TE', 1)
        power_tag = generate_sensor_tag(equipment_tag, 'W', 1) # Assuming W is power here
        pressure_tag = generate_sensor_tag(equipment_tag, 'P', 1)
        content.append(f"{temp_tag}: {random.randint(1450, 1600)} °C")
        content.append(f"{power_tag}: {random.uniform(20,30):.1f} MW")
        content.append(f"{pressure_tag}: {random.uniform(95, 115):.0f} kPa")
    elif equipment_type == "Crusher":
        load_tag = generate_sensor_tag(equipment_tag, 'W', 1) # Assuming W is load
        power_tag = generate_sensor_tag(equipment_tag, 'W', 2) # Assuming W is power
        vibration_tag = generate_sensor_tag(equipment_tag, 'V', 1) # Vibration
        content.append(f"{load_tag}: {random.uniform(80, 120):.1f} t/h")
        content.append(f"{power_tag}: {random.uniform(15, 25):.1f} MW")
        content.append(f"{vibration_tag}: {random.uniform(2, 8):.1f} mm/s")
    elif equipment_type == "Refinery":
        temp_tag = generate_sensor_tag(equipment_tag, 'TE', 1)
        o2_flow_tag = generate_sensor_tag(equipment_tag, 'F', 1) # O2 Flow
        content.append(f"{temp_tag}: {random.randint(1550, 1650)} °C")
        content.append(f"{o2_flow_tag}: {random.uniform(450, 550):.0f} Nm³/h")
    elif equipment_type == "Boiler":
        temp_tag = generate_sensor_tag(equipment_tag, 'TE', 1)
        pressure_tag = generate_sensor_tag(equipment_tag, 'P', 1)
        level_tag = generate_sensor_tag(equipment_tag, 'L', 1) # Level
        content.append(f"{temp_tag}: {random.randint(380, 420)} °C")
        content.append(f"{pressure_tag}: {random.uniform(100, 120):.0f} kPa")
        content.append(f"{level_tag}: {random.uniform(60, 90):.0f} %")


    content.append(f"Status: {status}")
    content.append("--- Alarms ---")
    if status == "ALARM":
        active_alarm_id = get_random_id("alarms", plant_code) or f"ALM-{plant_code}-{generate_uuid()}" # Get a real one if possible
        content.append(f"Active Alarm ID: {active_alarm_id.split('-')[-1]}") # Show short ID
        content.append(f"Alarm Type: {random.choice(['Temp High', 'Pressure Low', 'Overload', 'Vibration High'])}")
    elif status == "MAINTENANCE":
         mwo_id = get_random_id("maintenance", plant_code) or f"MWO-{plant_code}-{generate_uuid()}"
         content.append(f"Under Maintenance (MWO: {mwo_id.split('-')[-1]})")
    else:
        content.append("Active Alarms: None")

    filepath = create_diagram_pil(filename_png, title, content, plant_code, category)
    add_metadata(filepath, plant_code, category, "PNG", f"Simulated HMI snapshot for {equipment_tag}.", timestamp, [equipment_tag], linked_ids={"ActiveAlarmID": [active_alarm_id] if active_alarm_id else [], "MWO_ID": [mwo_id] if status == "MAINTENANCE" else []})


def generate_rich_media(plant_site_info):
    """Generates P&IDs, HMIs, other diagrams using standard tags."""
    plant_code = plant_site_info["code"]
    category = "Diagrams_Visuals"
    data_dir = os.path.join(base_dir, plant_code, category)
    ensure_dir(data_dir)
    print(f"  Generating Rich Media for {plant_code}...")
    current_date_str = datetime.now().strftime('%Y%m%d')
    equipment_list = current_plant_data.get("equipment_list", [])

    # Simplified P&ID for Raw Material Handling
    pid_filename = f"{plant_code}_PID_RawMat_{current_date_str}.pdf"
    generate_simplified_pid(pid_filename, plant_code)

    # HMI Snapshots for some major equipment
    hmi_candidates = [e for e in equipment_list if e['type'] in ['EAF', 'Crusher', 'Refinery', 'Boiler']]
    num_hmis = min(len(hmi_candidates), 3) # Generate for up to 3 pieces of equipment
    if hmi_candidates:
        for equip in random.sample(hmi_candidates, num_hmis):
            hmi_filename = f"{equip['tag']}_HMI_Snapshot_{datetime.now().strftime('%Y%m%d%H%M%S')}.png"
            generate_hmi_snapshot(hmi_filename, plant_code, equip['tag'], equip['type'])
    else:
        print(f"    WARN: No equipment found for HMI generation in {plant_code}")


    # Safety Checklist (PNG) - Less dynamic, more static content
    safety_filename = f"{plant_code}_Safety_Checklist_{current_date_str}.png"
    random_area_code = get_random_area_code(plant_code)
    random_area_name = plant_site_info["areas"].get(random_area_code, "General Area")
    mwo_link_id = get_random_id("maintenance", plant_code) or f"MWO-{plant_code}-{generate_uuid()}"
    gas_sensor_tag = generate_sensor_tag(f"{plant_code}-{random_area_code}-AREA", 'A', 1) # Generic area gas sensor

    safety_content = [
        f"DAILY SAFETY INSPECTION - {plant_code}",
        f"Date: {current_date_str}",
        f"Inspector: {random.choice([r for r in PERSONNEL_ROLES if 'Supervisor' in r or 'Safety' in r])}",
        f"Area: {random_area_name} ({random_area_code})",
        "---",
        "[X] Emergency exits clear",
        "[X] Fire extinguishers checked (Tag: FE-{random_area_code}-01)",
        f"[ ] PPE stock levels adequate (Action: Restock gloves - Ref: {mwo_link_id.split('-')[-1]})", # Link to MWO
        "[X] Walkways clear of obstructions",
        f"[X] Gas detection system online (Tag: {gas_sensor_tag})"
    ]
    safety_filepath = create_diagram_pil(safety_filename, "Safety Checklist", safety_content, plant_code, category)
    add_metadata(safety_filepath, plant_code, category, "PNG", f"Daily safety inspection checklist for area {random_area_code}.", current_date_str, linked_ids={"MWO_ID": [mwo_link_id]})

    print(f"    Rich media generation for {plant_code} complete.")

# --- 7. Safety Incident Reports ---
def generate_safety_incident_reports(plant_site_info, num_days=90):
    """Generates safety incident reports (DOCX/PDF)."""
    plant_code = plant_site_info["code"]
    plant_name = plant_site_info["name"]
    category = "Safety_Data"
    data_dir = os.path.join(base_dir, plant_code, category)
    ensure_dir(data_dir)
    print(f"  Generating Safety Incident Reports for {plant_code}...")

    end_date = datetime.now()
    start_date = end_date - timedelta(days=num_days)
    num_incidents = 15 + random.randint(-5, 10) # Fewer incidents than alarms/maint
    generated_ids["incidents"] = [] # Reset

    for i in range(num_incidents):
        incident_id = f"INC-{plant_code}-{generate_uuid()}"
        generated_ids["incidents"].append(incident_id)
        incident_date = start_date + timedelta(days=random.randint(0, num_days-1), hours=random.randint(6, 20)) # Occur during work hours
        incident_type = random.choice(INCIDENT_TYPES)
        severity = random.choice(SEVERITY_LEVELS)
        if incident_type == "Near Miss": severity = "Minor"
        elif "Injury" in incident_type: severity = random.choice(["Minor", "Moderate", "Serious"])
        else: severity = random.choice(["Minor", "Moderate", "Serious", "Major"])

        location_area_code = get_random_area_code(plant_code)
        location_desc = f"{plant_site_info['areas'].get(location_area_code)} Area"
        equipment_involved = get_random_equipment_tag(plant_code) if random.random() < 0.7 else "N/A"
        personnel_involved = [random.choice(PERSONNEL_ROLES) for _ in range(random.randint(1, 3))]
        reporter = random.choice([r for r in PERSONNEL_ROLES if "Supervisor" in r or "Safety" in r])

        # Describe the incident based on type and hazard
        hazard = random.choice(COMMON_HAZARDS)
        description = f"Incident Type: {incident_type}. Severity: {severity}. Hazard: {hazard}. "
        if incident_type == "Near Miss":
            description += f"Personnel [{', '.join(personnel_involved)}] narrowly avoided {hazard} near {equipment_involved} in {location_desc}."
        elif "Injury" in incident_type:
            description += f"Personnel [{personnel_involved[0]}] sustained a {severity.lower()} injury due to {hazard} while operating/working near {equipment_involved}."
            if severity == "Serious": description += " Medical transport required."
        elif incident_type == "Property Damage":
            description += f"Damage occurred to {equipment_involved} due to {hazard}. Estimated cost: ${random.randint(500, 50000)}."
        elif incident_type == "Environmental Spill":
            description += f"Minor spill of {random.choice(['hydraulic oil', 'lubricant', 'process water'])} near {equipment_involved}. Contained and cleaned up."

        # Immediate Actions
        actions_taken = ["Area secured.", "First aid administered."] if "Injury" in incident_type else ["Area inspected."]
        if severity in ["Serious", "Major"] or incident_type == "Property Damage":
            actions_taken.append("Equipment shut down for inspection.")
            actions_taken.append(f"Maintenance requested (Ref MWO-{plant_code}-{generate_uuid()}).") # Link to potential MWO
        if severity in ["Serious", "Major"]:
             actions_taken.append(f"RCA initiated (Ref RCA-{plant_code}-{generate_uuid()}).") # Link to potential RCA

        # --- Generate DOCX Report ---
        doc = DocxDocument()
        doc.add_heading(f"Safety Incident Report - {plant_name}", level=1)
        doc.add_paragraph().add_run(f"Incident ID: {incident_id}").bold = True
        doc.add_paragraph(f"Date & Time: {incident_date.strftime('%Y-%m-%d %H:%M')}")
        doc.add_paragraph(f"Location: {location_desc} (Area: {location_area_code})")
        doc.add_paragraph(f"Equipment Involved: {equipment_involved}")
        doc.add_paragraph(f"Incident Type: {incident_type}")
        doc.add_paragraph(f"Severity: {severity}")
        doc.add_paragraph(f"Personnel Involved: {', '.join(personnel_involved)}")
        doc.add_paragraph(f"Reported By: {reporter}")
        doc.add_heading("Incident Description", level=2)
        doc.add_paragraph(description)
        doc.add_heading("Immediate Actions Taken", level=2)
        for action in actions_taken:
            doc.add_paragraph(action, style='List Bullet')
        doc.add_heading("Potential Contributing Factors (Initial)", level=2)
        doc.add_paragraph(f"- {random.choice(['Human error', 'Equipment malfunction', 'Procedure not followed', 'Inadequate PPE', 'Poor housekeeping'])}", style='List Bullet')
        if random.random() > 0.5:
             doc.add_paragraph(f"- {random.choice(['Lack of training', 'Communication issue', 'Design flaw', 'External factor'])}", style='List Bullet')

        doc_path = os.path.join(data_dir, f"{incident_id}_SafetyIncidentReport.docx")
        doc.save(doc_path)
        # Extract linked IDs from actions taken
        linked_mwo = next((action.split('MWO-')[-1].strip(').') for action in actions_taken if "MWO-" in action), None)
        linked_rca = next((action.split('RCA-')[-1].strip(').') for action in actions_taken if "RCA-" in action), None)
        linked_ids_dict = {}
        if linked_mwo: linked_ids_dict["MWO_ID"] = [f"MWO-{plant_code}-{linked_mwo}"]
        if linked_rca: linked_ids_dict["RCA_ID"] = [f"RCA-{plant_code}-{linked_rca}"]

        add_metadata(doc_path, plant_code, category, "DOCX", f"Safety Incident Report for {incident_type} ({severity}).", incident_date.strftime('%Y-%m-%d'), [equipment_involved] if equipment_involved != "N/A" else [], linked_ids_dict)

    print(f"    Safety Incident Report generation for {plant_code} complete.")

# --- 8. Standard Operating Procedures (SOPs) ---
def generate_sops(plant_site_info):
    """Generates diverse SOPs (PDF) for different roles and tasks."""
    plant_code = plant_site_info["code"]
    plant_name = plant_site_info["name"]
    category = "Reports_SOPs" # Keep with Prod Reports
    data_dir = os.path.join(base_dir, plant_code, category)
    ensure_dir(data_dir)
    print(f"  Generating Standard Operating Procedures (SOPs) for {plant_code}...")

    sop_templates = [
        {"role": "Furnace Operator", "task": "EAF Tapping", "area": "SMT", "equip_type": "EAF", "doc_id": "TAP-001",
         "steps": ["Verify furnace temp ({TE_SENSOR}) & pre-tap analysis.", "Ensure ladle is preheated and positioned.", "Wear full aluminized PPE.", "Clear tap-hole area.", "Open tap-hole using oxygen lance/drill.", "Monitor metal flow into ladle.", "Close tap-hole securely.", "Take ladle sample (Ref Lab SOP LSA-002)."],
         "hazards": ["Molten metal splash", "High temperature", "Radiant heat"]},
        {"role": "Maintenance Technician", "task": "Crusher Bearing Lubrication", "area": "RMH", "equip_type": "CRUSH", "doc_id": "LUB-003",
         "steps": ["Perform Lock-Out/Tag-Out (LOTO) on {EQUIP_TAG}.", "Clean grease fittings.", "Apply specified lubricant (Type XYZ) using grease gun.", "Check for leaks or excessive grease.", "Clean up any spills.", "Remove LOTO and test run."],
         "hazards": ["Moving machinery", "Pinch points", "Stored energy"]},
        {"role": "Lab Analyst", "task": "Metal Sample Preparation", "area": "LAB", "equip_type": "LAB-EQUIP", "doc_id": "LSA-002",
         "steps": ["Receive ladle sample (Ref Batch {BATCH_ID}).", "Allow sample to cool.", "Cut sample using abrasive saw (Wear face shield/gloves).", "Grind sample surface until smooth.", "Analyze sample using Spectrometer (Ref SOP LSA-005).", "Record results against Batch ID in LIMS."],
         "hazards": ["Sharp edges", "Hot metal", "Dust inhalation", "Noise"]},
         {"role": "Refinery Operator", "task": "Starting Oxygen Lance", "area": "REF", "equip_type": "REF", "doc_id": "OXY-001",
         "steps": ["Confirm correct lance position above bath.", "Verify cooling water flow to lance.", "Check oxygen pressure and flow setpoints.", "Initiate oxygen flow gradually.", "Monitor lance condition and bath reaction.", "Adjust flow as per process requirements (Target Si: {SI_TARGET}%)"],
         "hazards": ["Molten metal ejection", "High pressure gas", "Lance failure"]},
         {"role": "Electrician", "task": "Electrode Slipping", "area": "SMT", "equip_type": "EAF", "doc_id": "ELE-002",
         "steps": ["Coordinate with Furnace Operator.", "Ensure power is off to the electrode holder.", "Verify hydraulic pressure for clamps.", "Carefully release electrode clamps.", "Allow electrode to slip to desired position using winch/crane.", "Re-tighten clamps to specified torque.", "Confirm electrode position and electrical contact."],
         "hazards": ["Electrical shock", "Falling objects (electrode)", "Hydraulic pressure release"]},
    ]

    num_sops_to_generate = 5 + random.randint(0, 3)
    equipment_list = current_plant_data.get("equipment_list", [])

    for i in range(num_sops_to_generate):
        template = random.choice(sop_templates)
        role = template["role"]
        task = template["task"]
        area = template["area"]
        equip_type = template["equip_type"]
        doc_id_suffix = template["doc_id"]

        # Find a relevant equipment tag
        relevant_equip = [e for e in equipment_list if e['type'] == equip_type and area in e['tag']]
        equip_tag = random.choice([e['tag'] for e in relevant_equip]) if relevant_equip else generate_equipment_tag(plant_code, area, equip_type, 1)

        doc_id = f"SOP-{plant_code}-{area}-{doc_id_suffix}"
        revision = f"{random.randint(1, 5)}.{random.randint(0, 9)}"
        sop_date = (datetime.now() - timedelta(days=random.randint(30, 365))).strftime("%Y-%m-%d")

        pdf_filename = os.path.join(data_dir, f"{doc_id}_Rev{revision}_{task.replace(' ','')}.pdf")
        doc = SimpleDocTemplate(pdf_filename, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        story.append(Paragraph(f"Standard Operating Procedure - {plant_name}", styles['h1']))
        story.append(Paragraph(f"Task: {task} ({equip_tag if equip_type != 'LAB-EQUIP' else 'Lab Area'})", styles['h2']))
        story.append(Paragraph(f"Document ID: {doc_id} Rev {revision}", styles['Normal']))
        story.append(Paragraph(f"Effective Date: {sop_date}", styles['Normal']))
        story.append(Paragraph(f"Applies To: {role}", styles['Normal']))
        story.append(Spacer(1, 12))

        story.append(Paragraph("1. Purpose:", styles['h3']))
        story.append(Paragraph(f"To define the standard procedure for safely performing '{task}'.", styles['Normal']))
        story.append(Paragraph("2. Scope:", styles['h3']))
        story.append(Paragraph(f"Applies to all {role}s working in the {plant_site_info['areas'].get(area)} ({area}).", styles['Normal']))
        story.append(Paragraph("3. Safety Precautions:", styles['h3']))
        # Change: Create a ListBullet style
        list_bullet_style = copy.deepcopy(styles['Normal']) # Ensure new instance
        list_bullet_style.spaceBefore = 6 # Add a bit of space
        list_bullet_style.leftIndent = 24 # Indent for bullet point
        story.append(Paragraph(f"Required PPE: {random.choice(['Standard PPE', 'Heat resistant PPE', 'Aluminized Suit', 'Chemical resistant gloves'])}", style=list_bullet_style))
        for hazard in template["hazards"]:
            story.append(Paragraph(f"Hazard: {hazard}", style=list_bullet_style))
        if "LOTO" in template["steps"][0]:
             story.append(Paragraph(f"Lock-Out/Tag-Out required on {equip_tag}.", style=list_bullet_style))

        story.append(Paragraph("4. Procedure Steps:", styles['h3']))
        step_num = 1
        for step in template["steps"]:
            # Replace placeholders
            step_text = step.replace("{EQUIP_TAG}", equip_tag)
            step_text = step_text.replace("{TE_SENSOR}", generate_sensor_tag(equip_tag, 'TE', 1))
            step_text = step_text.replace("{BATCH_ID}", get_random_id("batches", plant_code) or "Current Batch")
            step_text = step_text.replace("{SI_TARGET}", str(random.uniform(0.5, 1.5))) # Example target
            story.append(Paragraph(f"4.{step_num}. {step_text}", styles['Normal']))
            step_num += 1

        story.append(Spacer(1, 12))
        story.append(Paragraph("5. Related Documents:", styles['h3']))
        story.append(Paragraph(f"- General Plant Safety Rules (SAF-GEN-001)", style=list_bullet_style))
        if "LOTO" in template["steps"][0]: story.append(Paragraph(f"- LOTO Procedure (SAF-LOTO-001)", style=list_bullet_style))
        # Add cross-references if applicable
        if "LSA-002" in " ".join(template["steps"]): story.append(Paragraph(f"- Lab Sample Prep SOP (SOP-{plant_code}-LAB-LSA-002)", style=list_bullet_style))
        if "LSA-005" in " ".join(template["steps"]): story.append(Paragraph(f"- Spectrometer Operation SOP (SOP-{plant_code}-LAB-LSA-005)", style=list_bullet_style))


        try:
            doc.build(story)
            add_metadata(pdf_filename, plant_code, category, "PDF", f"SOP for {role} performing {task}.", sop_date, [equip_tag] if equip_type != 'LAB-EQUIP' else [], linked_ids={"BatchID": [bid for bid in [get_random_id("batches", plant_code)] if bid]})
        except Exception as e:
            print(f"    ERROR generating PDF SOP {pdf_filename}: {e}")

    print(f"    SOP generation for {plant_code} complete.")

# --- 9. Root Cause Analysis (RCA) Reports ---
def generate_rca_reports(plant_site_info, num_days=90):
    """Generates RCA reports (DOCX/PPTX) for significant incidents or downtime."""
    plant_code = plant_site_info["code"]
    plant_name = plant_site_info["name"]
    category = "RCA_Reports"
    data_dir = os.path.join(base_dir, plant_code, category)
    ensure_dir(data_dir)
    print(f"  Generating Root Cause Analysis (RCA) Reports for {plant_code}...")

    end_date = datetime.now()
    start_date = end_date - timedelta(days=num_days)
    generated_ids["rcas"] = [] # Reset

    # Identify potential triggers for RCA: Serious/Major Incidents or Long/Costly Downtime/Maintenance
    serious_incidents = [i for i in generated_ids.get("incidents", []) if ("Serious" in i or "Major" in i) and f"-{plant_code}-" in i and random.random() < 0.8]
    major_downtime = [d for d in current_plant_data.get("downtime_data", []) if d["DurationHours"] > 24 and random.random() < 0.6]
    major_maintenance = [m for m in current_plant_data.get("maintenance_data", []) if (m["Cost"] > 5000 or m["MaintenanceType"] == "Breakdown") and random.random() < 0.5]

    rca_triggers = serious_incidents + [d["DowntimeID"] for d in major_downtime] + [m["MWO_ID"] for m in major_maintenance]
    num_rcas = min(len(rca_triggers), 5 + random.randint(0, 5)) # Limit number of RCAs

    if not rca_triggers:
        print(f"    No significant events found to trigger RCA for {plant_code}.")
        return

    selected_triggers = random.sample(rca_triggers, num_rcas)

    for trigger_id in selected_triggers:
        rca_id = f"RCA-{plant_code}-{generate_uuid()}"
        generated_ids["rcas"].append(rca_id)
        rca_date = (end_date - timedelta(days=random.randint(0, num_days // 2))).strftime("%Y-%m-%d") # RCA happens after event
        rca_method = random.choice(RCA_METHODS)
        team_lead = random.choice([r for r in PERSONNEL_ROLES if "Supervisor" in r or "Safety" in r])
        team_members = random.sample([r for r in PERSONNEL_ROLES if r != team_lead], k=random.randint(2, 4))

        # --- Find details of the trigger event ---
        event_desc = f"Analysis of Event {trigger_id}"
        event_date = (end_date - timedelta(days=random.randint(num_days // 2, num_days))).strftime("%Y-%m-%d")
        equipment_tag = get_random_equipment_tag(plant_code) # Default
        linked_ids_dict = {}

        if trigger_id.startswith("INC-"):
            linked_ids_dict["IncidentID"] = [trigger_id]
            incident_info = next((meta for meta in metadata_list if trigger_id in meta["file_path"] and meta["category"] == "Safety_Data"), None)
            if incident_info:
                event_desc = incident_info["description"]
                event_date = incident_info["date_range"] # Single date for incident
                equipment_tag = incident_info["equipment_tags"][0] if incident_info["equipment_tags"] else equipment_tag
        elif trigger_id.startswith("DWT-"):
            linked_ids_dict["DowntimeID"] = [trigger_id]
            downtime_info = next((d for d in current_plant_data.get("downtime_data", []) if d["DowntimeID"] == trigger_id), None)
            if downtime_info:
                event_desc = f"Downtime Event: {downtime_info['Reason']} on {downtime_info['EquipmentTag']} ({downtime_info['DurationHours']}h)"
                event_date = downtime_info["StartTime"][:10]
                equipment_tag = downtime_info["EquipmentTag"]
                if downtime_info.get("TriggeringAlarmID"): linked_ids_dict["AlarmID"] = [downtime_info["TriggeringAlarmID"]]
                if downtime_info.get("RelatedMWO_ID"): linked_ids_dict["MWO_ID"] = [downtime_info["RelatedMWO_ID"]]
        elif trigger_id.startswith("MWO-"):
            linked_ids_dict["MWO_ID"] = [trigger_id]
            maint_info = next((m for m in current_plant_data.get("maintenance_data", []) if m["MWO_ID"] == trigger_id), None)
            if maint_info:
                event_desc = f"Maintenance Event: {maint_info['MaintenanceType']} on {maint_info['EquipmentTag']} (Cost: ${maint_info['Cost']:.0f})"
                event_date = maint_info["Date"]
                equipment_tag = maint_info["EquipmentTag"]
                if maint_info.get("TriggeringAlarmID"): linked_ids_dict["AlarmID"] = [maint_info["TriggeringAlarmID"]]
                if maint_info.get("RelatedDowntimeID"): linked_ids_dict["DowntimeID"] = [maint_info["RelatedDowntimeID"]]

        # --- Generate RCA Content ---
        # Simplified RCA generation
        problem_statement = f"Determine the root cause of: {event_desc}"
        # 5 Whys Example
        whys = [
            f"Why did {equipment_tag} fail/incident occur? - {random.choice(['Component failure', 'Incorrect operation', 'External factor'])}",
            f"Why did the component fail/operation error occur? - {random.choice(['Lack of maintenance', 'Inadequate training', 'Poor design'])}",
            f"Why was maintenance/training inadequate? - {random.choice(['Resource constraints', 'Procedure not updated', 'Complacency'])}",
            f"Why were resources constrained/procedure outdated? - {random.choice(['Budget cuts', 'Management oversight', 'Process change ignored'])}",
            f"Why was there oversight/process ignored? - {random.choice(['Lack of communication', 'Insufficient review process', 'Systemic issue'])}"
        ]
        root_cause = whys[-1].split('? - ')[-1]
        contributing_factors = [whys[i].split('? - ')[-1] for i in range(len(whys)-1)] + [random.choice(COMMON_HAZARDS)]
        corrective_actions = [
            f"Update SOP-{plant_code}-{equipment_tag.split('-')[1]}-{random.choice(['TAP','LUB','ELE'])}-XXX",
            f"Conduct refresher training for {random.choice(PERSONNEL_ROLES)}s",
            f"Implement enhanced pre-use inspection checklist for {equipment_tag}",
            f"Review maintenance schedule for {equipment_tag.split('-')[2]}",
            f"Assign responsibility: {random.choice(team_members)} by {(datetime.strptime(rca_date, '%Y-%m-%d') + timedelta(days=30)).strftime('%Y-%m-%d')}"
        ]

        # --- Generate DOCX RCA Report ---
        doc = DocxDocument()
        doc.add_heading(f"Root Cause Analysis Report - {plant_name}", level=1)
        doc.add_paragraph().add_run(f"RCA ID: {rca_id}").bold = True
        doc.add_paragraph(f"Date of Analysis: {rca_date}")
        doc.add_paragraph(f"Triggering Event ID: {trigger_id}")
        doc.add_paragraph(f"Event Date: {event_date}")
        doc.add_paragraph(f"Equipment/Area: {equipment_tag}")
        doc.add_paragraph(f"RCA Method: {rca_method}")
        doc.add_paragraph(f"Team Lead: {team_lead}")
        doc.add_paragraph(f"Team Members: {', '.join(team_members)}")

        doc.add_heading("1. Problem Statement", level=2)
        doc.add_paragraph(problem_statement)

        doc.add_heading("2. Analysis (Example: 5 Whys)", level=2)
        if rca_method == "5 Whys":
            for why in whys:
                doc.add_paragraph(why, style='List Bullet')
        else: # Simplified text for other methods
            doc.add_paragraph(f"Analysis using {rca_method} identified several potential causal chains.")
            doc.add_paragraph("Key factors considered: Human Factors, Equipment Condition, Procedures, Environment.")

        doc.add_heading("3. Contributing Factors", level=2)
        for factor in contributing_factors:
            doc.add_paragraph(factor, style='List Bullet')

        doc.add_heading("4. Identified Root Cause(s)", level=2)
        doc.add_paragraph(root_cause)

        doc.add_heading("5. Recommended Corrective Actions", level=2)
        for action in corrective_actions:
            doc.add_paragraph(action, style='List Bullet')

        doc_path = os.path.join(data_dir, f"{rca_id}_RCAReport.docx")
        doc.save(doc_path)
        add_metadata(doc_path, plant_code, category, "DOCX", f"RCA Report for event {trigger_id} using {rca_method}.", rca_date, [equipment_tag] if equipment_tag != "N/A" else [], linked_ids_dict)

        # --- Generate PPTX Summary ---
        try:
            prs = Presentation()
            # Title Slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = f"RCA Summary: Event {trigger_id}"
            subtitle.text = f"{plant_name} ({plant_code})\nDate: {rca_date}\nRCA ID: {rca_id}"

            # Content Slide Layout
            bullet_slide_layout = prs.slide_layouts[1]

            # Problem Description Slide
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = "Event Description"
            tf = body_shape.text_frame
            tf.text = f"Event ID: {trigger_id}"
            p = tf.add_paragraph()
            p.text = f"Date: {event_date}"
            p = tf.add_paragraph()
            p.text = f"Equipment: {equipment_tag}"
            p = tf.add_paragraph()
            p.text = f"Summary: {event_desc[:150]}..." # Keep it brief

            # Root Cause Slide
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = "Identified Root Cause & Contributing Factors"
            tf = body_shape.text_frame
            tf.text = f"Root Cause: {root_cause}"
            p = tf.add_paragraph()
            p.text = "Contributing Factors:"
            p.level = 1
            for factor in contributing_factors[:3]: # Max 3 factors for slide
                p = tf.add_paragraph()
                p.text = factor
                p.level = 2

            # Corrective Actions Slide
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = "Corrective Actions"
            tf = body_shape.text_frame
            for action in corrective_actions[:4]: # Max 4 actions
                 p = tf.add_paragraph()
                 p.text = action

            pptx_path = os.path.join(data_dir, f"{rca_id}_RCASummary.pptx")
            prs.save(pptx_path)
            add_metadata(pptx_path, plant_code, category, "PPTX", f"RCA Summary presentation for event {trigger_id}.", rca_date, [equipment_tag] if equipment_tag != "N/A" else [], linked_ids_dict)

        except Exception as e:
            print(f"    ERROR generating PPTX RCA Summary {rca_id}: {e}")


    print(f"    RCA Report generation for {plant_code} complete.")


# --- 10. Metadata Generation ---
def generate_metadata_file():
    """Generates a JSON file containing metadata for all generated files."""
    metadata_path = os.path.join(base_dir, "metadata_inventory.json")
    print(f"\nGenerating metadata inventory file: {metadata_path}")
    try:
        # Sort metadata for consistency by file path
        sorted_metadata = sorted(metadata_list, key=lambda x: x['file_path'])
        with open(metadata_path, 'w') as f:
            json.dump(sorted_metadata, f, indent=2)
        print("Metadata inventory generated successfully.")
    except Exception as e:
        print(f"ERROR generating metadata file: {e}")

# --- Main Execution ---
if __name__ == "__main__":
    # Global metadata list accumulates across plants
    metadata_list = []

    for plant_info in PLANT_SITES:
        plant_code = plant_info["code"]
        print(f"\nProcessing Plant: {plant_info['name']} ({plant_code})")

        # --- Reset plant-specific globals ---
        generated_ids = {
            "alarms": [], "maintenance": [], "downtime": [], "batches": [],
            "incidents": [], "rcas": []
        }
        current_plant_data = {} # Holds lists like equipment_list, alarm_data etc. for this plant run

        # Create plant-specific base directory
        plant_base_dir = os.path.join(base_dir, plant_code)
        ensure_dir(plant_base_dir)

        # --- Generate data sequentially for potential linking ---
        # Define history duration
        num_days_history = 90
        num_months_history = num_days_history // 30

        # 1. Core Equipment & Sensor Data (shorter history for sensor files)
        generate_sensor_data(plant_info, days=max(7, num_days_history // 12))

        # 2. Operational Logs (longer history, uses equipment list)
        generate_operational_logs(plant_info, num_days=num_days_history)

        # 3. Lab/Quality Data (linked to batches)
        generate_lab_reports(plant_info, num_days=num_days_history)

        # 4. Production Report (summarizes recent month)
        generate_production_report(plant_info, num_days=30)

        # 5. Environmental Reports (monthly for recent history)
        generate_environmental_reports(plant_info, num_months=num_months_history)

        # 6. Rich Media (Diagrams, HMIs - uses equipment list)
        generate_rich_media(plant_info)

        # 7. Safety Incidents (occur over history)
        generate_safety_incident_reports(plant_info, num_days=num_days_history)

        # 8. SOPs (Static, but generated per plant)
        generate_sops(plant_info)

        # 9. RCAs (Triggered by significant events in history)
        generate_rca_reports(plant_info, num_days=num_days_history)

        # --- End of plant processing ---
        print(f"Finished processing {plant_code}.")


    # 10. Generate the final metadata file covering all plants
    generate_metadata_file()

    print(f"\n\n---\nAll datasets successfully created in: {base_dir}")
    for plant_info in PLANT_SITES:
        print(f"  - {plant_info['code']} data in: {os.path.join(base_dir, plant_info['code'])}")
    print(f"  - Metadata inventory: {os.path.join(base_dir, 'metadata_inventory.json')}")